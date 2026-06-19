#!/usr/bin/env python3
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# COLORES DEL PROYECTO
DARK_BG = RGBColor(9, 13, 22)          # Slate oscuro (#090D16)
LIGHT_BG = RGBColor(248, 250, 252)     # Slate claro (#F8FAFC)
BRAND_RED = RGBColor(220, 38, 38)      # Rojo institucional (#DC2626)
BRAND_EMERALD = RGBColor(5, 150, 105)  # Verde clínico (#059669)
TEXT_MAIN = RGBColor(15, 23, 42)       # Slate 900 (#0F172A)
TEXT_MUTED = RGBColor(71, 85, 105)     # Slate 600 (#475569)
TEXT_WHITE = RGBColor(255, 255, 255)   # Blanco
BORDER_COLOR = RGBColor(226, 232, 240) # Gris claro (#E2E8F0)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, category, title):
    # Categoría (Texto pequeño en verde clínico)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "●  " + category.upper()
    p.font.name = 'Montserrat'
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = BRAND_EMERALD
    
    # Título (Texto grande en gris oscuro)
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.name = 'Montserrat'
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MAIN

def add_slide_footer(slide):
    # Línea sutil horizontal
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.015))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER_COLOR
    shape.line.fill.background()
    
    # Marca izquierda
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(6.0), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "🏥  Cooperadora Hospital Emilio Ferreyra"
    p.font.name = 'Inter'
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = BRAND_EMERALD
    
    # Curso / Entrega derecha
    txBox2 = slide.shapes.add_textbox(Inches(8.5), Inches(6.9), Inches(4.0), Inches(0.3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    p2.text = "Metodología de Sistemas 2 · TFI"
    p2.font.name = 'Inter'
    p2.font.size = Pt(9)
    p2.font.color.rgb = TEXT_MUTED

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # DIAPOSITIVA 1: PORTADA (Tema Oscuro Premium)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, DARK_BG)
    
    # Universidad / Categoría
    tx = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(0.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "UNIVERSIDAD TECNOLÓGICA NACIONAL (UTN)  ·  EXTENSIÓN ÁULICA NECOCHEA"
    p.font.name = 'Montserrat'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BRAND_EMERALD
    
    # Materia e Info
    tx_tfi = slide1.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(0.4))
    tf_tfi = tx_tfi.text_frame
    tf_tfi.word_wrap = True
    p_tfi = tf_tfi.paragraphs[0]
    p_tfi.text = "Metodología de Sistemas 2  |  Trabajo Final Integrador"
    p_tfi.font.name = 'Montserrat'
    p_tfi.font.size = Pt(14)
    p_tfi.font.bold = True
    p_tfi.font.color.rgb = BRAND_RED
    
    # Título Principal
    tx_title = slide1.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.333), Inches(2.2))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = "Portal Web de la Asociación Cooperadora del\nHospital Municipal \"Dr. Emilio Ferreyra\""
    p_title.font.name = 'Montserrat'
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    
    # Línea decorativa
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.5), Inches(11.333), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_EMERALD
    line.line.fill.background()
    
    # Miembros del Equipo
    tx_team = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(6.5), Inches(2.0))
    tf_team = tx_team.text_frame
    tf_team.word_wrap = True
    tf_team.margin_left = tf_team.margin_top = 0
    p_t = tf_team.paragraphs[0]
    p_t.text = "EQUIPO DE TRABAJO Y ROLES:"
    p_t.font.name = 'Montserrat'
    p_t.font.size = Pt(10)
    p_t.font.bold = True
    p_t.font.color.rgb = BRAND_EMERALD
    p_t.space_after = Pt(8)
    
    members = [
        ("Aramis Prieto", "Scrum Master & Lead Backend Developer"),
        ("Kevin Nielsen", "Security & Backend Developer"),
        ("Thiago Masson", "Full Stack Developer & Branding Specialist"),
        ("Santiago Ialungo", "UI/UX & Lead Frontend Developer")
    ]
    for name, role in members:
        p_m = tf_team.add_paragraph()
        p_m.text = f"•  {name}  -  "
        p_m.font.name = 'Inter'
        p_m.font.size = Pt(10)
        p_m.font.bold = True
        p_m.font.color.rgb = TEXT_WHITE
        
        # Añadir el rol con estilo atenuado
        run = p_m.add_run()
        run.text = role
        run.font.name = 'Inter'
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED
    
    # Metadata Entrega (Derecha)
    tx_meta = slide1.shapes.add_textbox(Inches(8.5), Inches(4.8), Inches(3.8), Inches(2.0))
    tf_meta = tx_meta.text_frame
    tf_meta.word_wrap = True
    tf_meta.margin_top = 0
    
    p_meta = tf_meta.paragraphs[0]
    p_meta.alignment = PP_ALIGN.RIGHT
    p_meta.text = "DOCENTE:"
    p_meta.font.name = 'Montserrat'
    p_meta.font.size = Pt(10)
    p_meta.font.bold = True
    p_meta.font.color.rgb = BRAND_EMERALD
    p_meta.space_after = Pt(4)
    
    p_doc = tf_meta.add_paragraph()
    p_doc.alignment = PP_ALIGN.RIGHT
    p_doc.text = "Daniel Moreno"
    p_doc.font.name = 'Inter'
    p_doc.font.size = Pt(12)
    p_doc.font.bold = True
    p_doc.font.color.rgb = TEXT_WHITE
    p_doc.space_after = Pt(14)
    
    p_date_lbl = tf_meta.add_paragraph()
    p_date_lbl.alignment = PP_ALIGN.RIGHT
    p_date_lbl.text = "FECHA DE ENTREGA:"
    p_date_lbl.font.name = 'Montserrat'
    p_date_lbl.font.size = Pt(10)
    p_date_lbl.font.bold = True
    p_date_lbl.font.color.rgb = BRAND_EMERALD
    p_date_lbl.space_after = Pt(4)
    
    p_date = tf_meta.add_paragraph()
    p_date.alignment = PP_ALIGN.RIGHT
    p_date.text = "Junio 2026  ·  Necochea"
    p_date.font.name = 'Inter'
    p_date.font.size = Pt(11)
    p_date.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # DIAPOSITIVA 2: ALCANCE Y REQUERIMIENTOS
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, LIGHT_BG)
    add_slide_header(slide2, "Requerimientos generales", "Descripción y Alcance del Proyecto")
    add_slide_footer(slide2)
    
    # Grid de 3 columnas x 2 filas
    modules = [
        ("01. Gestión de Socios", "Registro de asociados con DNI único, domicilio, teléfono y alta automatizada al Libro de Registro digital."),
        ("02. Módulo Transaccional", "Historial de pagos de cuotas sociales y control de membresía con pasarela de pago recurrente (Mercado Pago)."),
        ("03. Recaudación Transparente", "Visualización de campañas de recaudación activas para insumos/obras, con barras de progreso en tiempo real."),
        ("04. Auditoría de Donaciones", "Checkout interactivo para registrar transferencias manuales (carga de comprobantes) o donación en línea directa."),
        ("05. Panel Administrativo", "Gestión de operadores para aprobar transferencias, modificar campañas, y auditar el padrón social."),
        ("06. Módulo Informativo", "Portal de novedades y de obras concretadas para difundir hitos institucionales de la cooperadora.")
    ]
    
    col_width = Inches(3.64)
    row_height = Inches(2.1)
    start_x = Inches(0.8)
    start_y = Inches(1.9)
    gap_x = Inches(0.4)
    gap_y = Inches(0.4)
    
    for i, (title, desc) in enumerate(modules):
        col = i % 3
        row = i // 3
        x = start_x + col * (col_width + gap_x)
        y = start_y + row * (row_height + gap_y)
        
        # Tarjeta
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_width, row_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.08  # esquinas suaves
        
        # Texto dentro de la tarjeta
        tx_box = slide2.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), col_width - Inches(0.4), row_height - Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Montserrat'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BRAND_EMERALD
        p.space_after = Pt(6)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Inter'
        p2.font.size = Pt(10)
        p2.font.color.rgb = TEXT_MUTED
        p2.line_spacing = 1.25

    # ----------------------------------------------------
    # DIAPOSITIVA 3: ROLES Y EQUIPO
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, LIGHT_BG)
    add_slide_header(slide3, "Planificación y Gestión", "Roles y Responsabilidades Técnicas")
    add_slide_footer(slide3)
    
    team_details = [
        ("Aramis Prieto", "Scrum Master & Lead Backend", [
            "Facilitación de metodologías ágiles.",
            "Modelado híbrido de datos (SQL+NoSQL).",
            "Transaccionalidad concurrente mediante bloqueos de fila relacionales.",
            "Desarrollo de suite de pruebas de integración con Vitest."
        ]),
        ("Kevin Nielsen", "Security & Backend Developer", [
            "Autenticación segura basada en tokens JWT.",
            "Rate limiting por IP contra ataques DDoS.",
            "Validaciones estructuradas con expresiones regulares.",
            "Servicio SMTP automatizado para notificaciones por correo electrónico."
        ]),
        ("Santiago Ialungo", "UI/UX & Lead Frontend", [
            "Diseño de interfaz clínica (paleta médica, ECG).",
            "Desarrollo responsivo con Tailwind CSS.",
            "Paneles públicos y panel administrativo de control.",
            "Lenis scroll e implementaciones de UI Scroll-Spy."
        ]),
        ("Thiago Masson", "Full Stack & Branding", [
            "Gestión de monorrepo con pnpm workspaces.",
            "Integración de logotipos institucionales.",
            "Desarrollo del módulo de noticias documentales.",
            "Sanitización preventiva contra ataques XSS con DOMPurify."
        ])
    ]
    
    col_w = Inches(2.73)
    start_x = Inches(0.8)
    gap_x = Inches(0.27)
    box_h = Inches(4.5)
    
    for i, (name, role, points) in enumerate(team_details):
        x = start_x + i * (col_w + gap_x)
        y = Inches(1.9)
        
        # Tarjeta de rol
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.06
        
        # Borde decorativo superior
        top_accent = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, Inches(0.12))
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = BRAND_EMERALD if i%2 == 0 else BRAND_RED
        top_accent.line.fill.background()
        top_accent.adjustments[0] = 0.5
        
        # Texto
        tx_box = slide3.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), col_w - Inches(0.4), box_h - Inches(0.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Nombre
        p_name = tf.paragraphs[0]
        p_name.text = name
        p_name.font.name = 'Montserrat'
        p_name.font.size = Pt(14)
        p_name.font.bold = True
        p_name.font.color.rgb = TEXT_MAIN
        
        # Rol
        p_role = tf.add_paragraph()
        p_role.text = role
        p_role.font.name = 'Inter'
        p_role.font.size = Pt(8.5)
        p_role.font.bold = True
        p_role.font.color.rgb = BRAND_EMERALD if i%2 == 0 else BRAND_RED
        p_role.space_after = Pt(12)
        
        # Puntos
        for pt in points:
            p_pt = tf.add_paragraph()
            p_pt.text = "⚡  "
            p_pt.font.name = 'Inter'
            p_pt.font.size = Pt(9.5)
            p_pt.font.bold = True
            p_pt.font.color.rgb = BRAND_RED
            p_pt.space_after = Pt(6)
            
            run = p_pt.add_run()
            run.text = pt
            run.font.name = 'Inter'
            run.font.size = Pt(9)
            run.font.bold = False
            run.font.color.rgb = TEXT_MUTED
            run.line_spacing = 1.15

    # ----------------------------------------------------
    # DIAPOSITIVA 4: ARQUITECTURA DE DATOS HÍBRIDA
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, LIGHT_BG)
    add_slide_header(slide4, "Arquitectura de Software", "Modelo Cliente-Servidor y Persistencia Híbrida")
    add_slide_footer(slide4)
    
    # Columna 1: Arquitectura General (Esquema de Cajas)
    # SPA
    spa = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(4.5), Inches(1.2))
    spa.fill.solid()
    spa.fill.fore_color.rgb = RGBColor(255, 255, 255)
    spa.line.color.rgb = BORDER_COLOR
    spa.line.width = Pt(1.5)
    spa.adjustments[0] = 0.1
    
    tx_spa = slide4.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(4.1), Inches(0.8))
    tf_spa = tx_spa.text_frame
    tf_spa.word_wrap = True
    p_spa = tf_spa.paragraphs[0]
    p_spa.text = "CLIENTE (Frontend SPA)"
    p_spa.font.name = 'Montserrat'
    p_spa.font.size = Pt(13)
    p_spa.font.bold = True
    p_spa.font.color.rgb = TEXT_MAIN
    p_spa2 = tf_spa.add_paragraph()
    p_spa2.text = "React.js  ·  Vite  ·  Tailwind CSS"
    p_spa2.font.name = 'Inter'
    p_spa2.font.size = Pt(10)
    p_spa2.font.color.rgb = BRAND_EMERALD
    p_spa2.font.bold = True
    
    # Flecha
    arrow = slide4.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(4.5), Inches(0.5))
    tf_arr = arrow.text_frame
    p_arr = tf_arr.paragraphs[0]
    p_arr.alignment = PP_ALIGN.CENTER
    p_arr.text = "⇅  HTTPS / REST API (JWT)"
    p_arr.font.name = 'Inter'
    p_arr.font.size = Pt(11)
    p_arr.font.bold = True
    p_arr.font.color.rgb = TEXT_MUTED
    
    # API
    api = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.6), Inches(4.5), Inches(1.2))
    api.fill.solid()
    api.fill.fore_color.rgb = RGBColor(255, 255, 255)
    api.line.color.rgb = BRAND_RED
    api.line.width = Pt(1.5)
    api.adjustments[0] = 0.1
    
    tx_api = slide4.shapes.add_textbox(Inches(1.0), Inches(3.8), Inches(4.1), Inches(0.8))
    tf_api = tx_api.text_frame
    tf_api.word_wrap = True
    p_api = tf_api.paragraphs[0]
    p_api.text = "API SERVER (Backend)"
    p_api.font.name = 'Montserrat'
    p_api.font.size = Pt(13)
    p_api.font.bold = True
    p_api.font.color.rgb = TEXT_MAIN
    p_api2 = tf_api.add_paragraph()
    p_api2.text = "Node.js  ·  Express  ·  pnpm workspaces"
    p_api2.font.name = 'Inter'
    p_api2.font.size = Pt(10)
    p_api2.font.color.rgb = BRAND_RED
    p_api2.font.bold = True
    
    # Columna 2: Persistencia Híbrida
    # SQL
    sql_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(1.9), Inches(3.2), Inches(2.2))
    sql_box.fill.solid()
    sql_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    sql_box.line.color.rgb = BRAND_EMERALD
    sql_box.line.width = Pt(1.5)
    sql_box.adjustments[0] = 0.1
    
    tx_sql = slide4.shapes.add_textbox(Inches(6.0), Inches(2.1), Inches(2.8), Inches(1.8))
    tf_sql = tx_sql.text_frame
    tf_sql.word_wrap = True
    p_sql = tf_sql.paragraphs[0]
    p_sql.text = "SQL (Transaccional ACID)"
    p_sql.font.name = 'Montserrat'
    p_sql.font.size = Pt(12)
    p_sql.font.bold = True
    p_sql.font.color.rgb = TEXT_MAIN
    p_sql.space_after = Pt(8)
    
    sqls = ["• usuarios: claves bcryptjs", "• perfiles_socios: datos DNI", "• campanas_eco: balances"]
    for s in sqls:
        p_s = tf_sql.add_paragraph()
        p_s.text = s
        p_s.font.name = 'Inter'
        p_s.font.size = Pt(9)
        p_s.font.color.rgb = TEXT_MUTED
        p_s.space_after = Pt(4)
        
    # NoSQL
    nosql_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.3), Inches(1.9), Inches(3.2), Inches(2.2))
    nosql_box.fill.solid()
    nosql_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    nosql_box.line.color.rgb = BRAND_RED
    nosql_box.line.width = Pt(1.5)
    nosql_box.adjustments[0] = 0.1
    
    tx_nosql = slide4.shapes.add_textbox(Inches(9.5), Inches(2.1), Inches(2.8), Inches(1.8))
    tf_nosql = tx_nosql.text_frame
    tf_nosql.word_wrap = True
    p_nosql = tf_nosql.paragraphs[0]
    p_nosql.text = "NoSQL (Documental Flexible)"
    p_nosql.font.name = 'Montserrat'
    p_nosql.font.size = Pt(12)
    p_nosql.font.bold = True
    p_nosql.font.color.rgb = TEXT_MAIN
    p_nosql.space_after = Pt(8)
    
    nosqls = ["• noticias_actualidad: galerías", "• campanas_detalle: multimedia", "  enlazado por campana_id_ref"]
    for ns in nosqls:
        p_ns = tf_nosql.add_paragraph()
        p_ns.text = ns
        p_ns.font.name = 'Inter'
        p_ns.font.size = Pt(9)
        p_ns.font.color.rgb = TEXT_MUTED
        p_ns.space_after = Pt(4)
        
    # Data Mashup & Locks (Bottom)
    mash_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(4.3), Inches(6.7), Inches(2.1))
    mash_box.fill.solid()
    mash_box.fill.fore_color.rgb = RGBColor(241, 245, 249) # Slate-100 background
    mash_box.line.color.rgb = BRAND_EMERALD
    mash_box.line.width = Pt(1)
    mash_box.adjustments[0] = 0.08
    
    tx_mash = slide4.shapes.add_textbox(Inches(6.0), Inches(4.5), Inches(6.3), Inches(1.7))
    tf_mash = tx_mash.text_frame
    tf_mash.word_wrap = True
    tf_mash.margin_left = tf_mash.margin_top = 0
    p_m = tf_mash.paragraphs[0]
    p_m.text = "MECANISMOS AVANZADOS DE DATOS:"
    p_m.font.name = 'Montserrat'
    p_m.font.size = Pt(11)
    p_m.font.bold = True
    p_m.font.color.rgb = BRAND_EMERALD
    p_m.space_after = Pt(6)
    
    p_m2 = tf_mash.add_paragraph()
    p_m2.text = "⚡  Data Mashup Sincrónico: "
    p_m2.font.name = 'Inter'
    p_m2.font.size = Pt(10)
    p_m2.font.bold = True
    p_m2.font.color.rgb = BRAND_RED
    run_m2 = p_m2.add_run()
    run_m2.text = "Las consultas a SQL (finanzas) y NoSQL (multimedia) se ejecutan en paralelo con Promise.all y se unifican en un JSON plano enviado al cliente."
    run_m2.font.bold = False
    run_m2.font.color.rgb = TEXT_MUTED
    p_m2.space_after = Pt(6)
    
    p_m3 = tf_mash.add_paragraph()
    p_m3.text = "⚡  Evitación de Condiciones de Carrera: "
    p_m3.font.name = 'Inter'
    p_m3.font.size = Pt(10)
    p_m3.font.bold = True
    p_m3.font.color.rgb = BRAND_RED
    run_m3 = p_m3.add_run()
    run_m3.text = "Las donaciones concurrentes bloquean la fila financiera en SQL (SELECT ... FOR UPDATE) garantizando consistencia y transaccionalidad ACID."
    run_m3.font.bold = False
    run_m3.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # DIAPOSITIVA 5: OPTIMIZACIÓN Y SEGURIDAD (Hardening)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, LIGHT_BG)
    add_slide_header(slide5, "Rendimiento y Hardening", "Optimización y Seguridad de Producción")
    add_slide_footer(slide5)
    
    # 2 Columnas
    # Columna 1: Frontend
    tx_front_title = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.4))
    tf_ft = tx_front_title.text_frame
    p_ft = tf_ft.paragraphs[0]
    p_ft.text = "✓  FRONTEND, UX & SEO"
    p_ft.font.name = 'Montserrat'
    p_ft.font.size = Pt(16)
    p_ft.font.bold = True
    p_ft.font.color.rgb = BRAND_EMERALD
    
    tx_front = slide5.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.2))
    tf_f = tx_front.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_top = 0
    
    front_opts = [
        ("Code Splitting & Lazy Load", "React.lazy() y Suspense en App.jsx para cargar perezosamente paneles pesados de administración."),
        ("Velocidad Perceptual (Skeletons)", "Uso de CampaignSkeleton y NewsSkeleton para eliminar parpadeo visual molesto."),
        ("Accesibilidad y Semántica SEO", "Etiquetas semánticas HTML5, etiquetas Open Graph completas y precarga DNS de tipografías.")
    ]
    for i, (title, desc) in enumerate(front_opts):
        p_t = tf_f.add_paragraph() if i > 0 else tf_f.paragraphs[0]
        p_t.text = f"✦  {title}"
        p_t.font.name = 'Inter'
        p_t.font.size = Pt(11.5)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MAIN
        p_t.space_after = Pt(2)
        
        p_d = tf_f.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(12)
        p_d.line_spacing = 1.15

    # Columna 2: Backend
    tx_back_title = slide5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(0.4))
    tf_bt = tx_back_title.text_frame
    p_bt = tf_bt.paragraphs[0]
    p_bt.text = "🛡  BACKEND, SEGURIDAD & CACHÉ"
    p_bt.font.name = 'Montserrat'
    p_bt.font.size = Pt(16)
    p_bt.font.bold = True
    p_bt.font.color.rgb = BRAND_RED
    
    tx_back = slide5.shapes.add_textbox(Inches(6.8), Inches(2.3), Inches(5.6), Inches(4.2))
    tf_b = tx_back.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_top = 0
    
    back_opts = [
        ("Caché Invalidador Reactivo", "NodeCache en lecturas públicas (5m TTL), con limpieza automática mediante flushCache() en altas o modificaciones."),
        ("Rate Limiting por IP", "Configuración express-rate-limit con trust proxy activo para prevenir ataques sin bloquear IPs de balanceadores."),
        ("Helmet & Sanitización NoSQL", "Middleware helmet para cabeceras HTTP robustas y mongo-sanitize para saneo de queries maliciosas.")
    ]
    for i, (title, desc) in enumerate(back_opts):
        p_t = tf_b.add_paragraph() if i > 0 else tf_b.paragraphs[0]
        p_t.text = f"🔒  {title}"
        p_t.font.name = 'Inter'
        p_t.font.size = Pt(11.5)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_MAIN
        p_t.space_after = Pt(2)
        
        p_d = tf_b.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Inter'
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(12)
        p_d.line_spacing = 1.15

    # ----------------------------------------------------
    # DIAPOSITIVA 6: PROPUESTAS DE REFACTORIZACIÓN
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, LIGHT_BG)
    add_slide_header(slide6, "Refactorización y Arquitectura", "Propuestas de Refactorización mediante Patrones de Diseño")
    add_slide_footer(slide6)
    
    patterns = [
        ("Singleton Pattern", "CONEXIÓN DE BD ÚNICA", "Módulo: config/db.js\nAsegura que toda la aplicación de Express comparta una única instancia de conexión para los motores SQL y NoSQL.\n\nVentaja: Ahorro de memoria, previene la sobrecarga del pool de conexiones al evitar múltiples instancias redundantes."),
        ("Factory Method Pattern", "CREACIÓN DE USUARIOS", "Módulo: services/userService.js\nDelegación de la creación de instancias del modelo según rol (Socio o Admin), precargando permisos específicos.\n\nVentaja: Cumple con el principio Open/Closed al poder añadir nuevos roles fácilmente sin alterar controladores."),
        ("State Pattern + Async Worker", "CONCILIACIÓN AUTOMÁTICA", "Módulo: donacionController.js\nFlujo de transferencias manuales dividido en clases de estado (EstadoPendiente, EstadoAprobado, EstadoRechazado) con un Cron Job / Worker de conciliación automática en segundo plano.\n\nVentaja: Resuelve el cuello de botella manual.")
    ]
    
    col_width = Inches(3.64)
    start_x = Inches(0.8)
    gap_x = Inches(0.4)
    box_height = Inches(4.5)
    
    for i, (pat_name, pat_subtitle, pat_desc) in enumerate(patterns):
        x = start_x + i * (col_width + gap_x)
        y = Inches(1.9)
        
        # Tarjeta
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_width, box_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = BORDER_COLOR
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.06
        
        # Texto
        tx_box = slide6.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), col_width - Inches(0.4), box_height - Inches(0.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = pat_name
        p.font.name = 'Montserrat'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BRAND_EMERALD
        
        p_sub = tf.add_paragraph()
        p_sub.text = pat_subtitle
        p_sub.font.name = 'Montserrat'
        p_sub.font.size = Pt(9.5)
        p_sub.font.bold = True
        p_sub.font.color.rgb = BRAND_RED
        p_sub.space_after = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = pat_desc
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(9)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.line_spacing = 1.25

    # ----------------------------------------------------
    # GUARDAR PRESENTACIÓN
    # ----------------------------------------------------
    prs.save("/Users/aramisprieto/Documents/cooperadora-hospital1/presentacion_cooperadora.pptx")
    print("Presentación PPTX generada con éxito en: presentacion_cooperadora.pptx")

if __name__ == '__main__':
    main()
